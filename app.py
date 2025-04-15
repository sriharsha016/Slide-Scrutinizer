from flask import Flask, render_template, request
from werkzeug.utils import secure_filename
import os
import hashlib
from pptx import Presentation
from googleapiclient.discovery import build
import time
import json
from difflib import SequenceMatcher

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
DATA_FILE = 'data.json'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

API_KEY = "AIzaSyDLwnq6LZbOdywpAtqmjLDx0hJpUEmGEz0"
SEARCH_ENGINE_ID = "b36d0933710b745ad"

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Clean start: remove old data
if os.path.exists(DATA_FILE):
    os.remove(DATA_FILE)

for filename in os.listdir(UPLOAD_FOLDER):
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    try:
        if os.path.isfile(file_path):
            os.remove(file_path)
    except Exception as e:
        print(f"Error removing file {file_path}: {e}")

def get_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + ' '
        return text.strip()
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ''

def calculate_hash(text):
    return hashlib.sha256(text.encode('utf-8')).hexdigest()

def search_google_for_content(content):
    service = build("customsearch", "v1", developerKey=API_KEY)
    try:
        res = service.cse().list(q=content, cx=SEARCH_ENGINE_ID).execute()
        total_results = int(res.get("searchInformation", {}).get("totalResults", 0))
        percent = min(100, total_results * 10)
        return percent
    except Exception as e:
        print(f"Google search error: {e}")
        return 0

def calculate_similarity(text1, text2):
    """Calculate similarity percentage between two texts"""
    if not text1 or not text2:
        return 0
    matcher = SequenceMatcher(None, text1, text2)
    similarity = matcher.ratio() * 100
    return round(similarity, 2)

def load_data():
    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, 'r') as f:
            return json.load(f)
    return {
        "uploaded_files": [],
        "skipped_files": [],
        "hashes": {},
        "google_percentages": [],
        "non_duplicates": [],
        "duplicates": [],
        "similarities": [],
        "file_texts": {}
    }

def save_data(data):
    with open(DATA_FILE, 'w') as f:
        json.dump(data, f, indent=4)

@app.route('/')
def index():
    data = load_data()
    return render_template('index.html',
                           uploaded=data["uploaded_files"],
                           skipped=data["skipped_files"],
                           duplicates=data["duplicates"],
                           google_percentages=data["google_percentages"],
                           non_duplicates=data["non_duplicates"],
                           similarities=data["similarities"])

@app.route('/upload', methods=['POST'])
def upload():
    data = load_data()

    if 'files' not in request.files:
        return render_template('index.html', error="No files uploaded.", **data)

    files = request.files.getlist('files')
    for file in files:
        if file and file.filename.endswith('.pptx'):
            original_filename = secure_filename(file.filename)
            unique_filename = f"{int(time.time() * 1000)}_{original_filename}"
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            file.save(filepath)

            text = get_text_from_pptx(filepath)
            data["uploaded_files"].append(unique_filename)

            if not text:
                data["skipped_files"].append(unique_filename)
                continue

            # Store the text for similarity comparison
            data["file_texts"][unique_filename] = text
            
            file_hash = calculate_hash(text)

            duplicate_found = False
            for hash_value, existing_filename in data["hashes"].items():
                if file_hash == hash_value:
                    data["duplicates"].append((existing_filename, unique_filename))
                    duplicate_found = True
                    break

            if not duplicate_found:
                data["hashes"][file_hash] = unique_filename
                percentage = search_google_for_content(text[:300])
                data["google_percentages"].append({
                    "filename": unique_filename,
                    "percentage": percentage
                })
                if percentage < 90:
                    data["non_duplicates"].append(unique_filename)
    
    # Calculate similarity between all non-duplicate files
    # Clear previous similarities
    data["similarities"] = []
    
    # Get all files that need to be compared (including the new ones)
    files_to_compare = list(data["file_texts"].keys())
    
    # Compare each file with every other file
    for i in range(len(files_to_compare)):
        for j in range(i+1, len(files_to_compare)):
            file1 = files_to_compare[i]
            file2 = files_to_compare[j]
            
            # Skip if they are already marked as duplicates
            is_duplicate = False
            for dup in data["duplicates"]:
                if (file1 == dup[0] and file2 == dup[1]) or (file1 == dup[1] and file2 == dup[0]):
                    is_duplicate = True
                    break
            
            if is_duplicate:
                continue
                
            # Calculate similarity
            similarity = calculate_similarity(data["file_texts"][file1], data["file_texts"][file2])
            
            # If similarity is above 50%, add to similarities list
            if similarity >= 50:
                data["similarities"].append({
                    "file1": file1,
                    "file2": file2,
                    "similarity": similarity
                })

    save_data(data)
    return render_template('index.html',
                           uploaded=data["uploaded_files"],
                           skipped=data["skipped_files"],
                           duplicates=data["duplicates"],
                           google_percentages=data["google_percentages"],
                           non_duplicates=data["non_duplicates"],
                           similarities=data["similarities"])

if __name__ == '__main__':
    app.run(debug=True)
